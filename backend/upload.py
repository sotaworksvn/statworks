"""POST /upload — Data Ingestion endpoint (F-01).

Updated for Phase 1.5: async R2 persistence + Supabase metadata.
"""

from __future__ import annotations

import asyncio
import logging
import uuid
from pathlib import PurePath

import pandas as pd
from fastapi import APIRouter, HTTPException, Request, UploadFile, File

from backend.auth.context import get_current_user_id
from backend.db import supabase as supa
from backend.models import ColumnMeta, PresignRequest, PresignResponse, UploadResponse
from backend.storage import r2
from backend.store import FileEntry, set_entry

logger = logging.getLogger(__name__)

router = APIRouter()

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
ALLOWED_PRIMARY = {".xlsx", ".csv"}
ALLOWED_CONTEXT = {".docx", ".pptx"}
ALLOWED_ALL = ALLOWED_PRIMARY | ALLOWED_CONTEXT
MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB
MAX_FILES_PER_UPLOAD = 5


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_extension(filename: str) -> str:
    """Return the lowercased file extension (e.g. '.xlsx')."""
    return PurePath(filename).suffix.lower()


def _parse_excel(content: bytes) -> pd.DataFrame:
    """Parse an .xlsx file into a DataFrame."""
    import io
    return pd.read_excel(io.BytesIO(content), engine="openpyxl")


def _parse_csv(content: bytes) -> pd.DataFrame:
    """Parse a .csv file into a DataFrame."""
    import io
    return pd.read_csv(io.BytesIO(content))


def _extract_docx_text(content: bytes) -> str:
    """Extract plain text from a .docx file."""
    import io
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx_text(content: bytes) -> str:
    """Extract plain text from a .pptx file."""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _detect_columns(df: pd.DataFrame) -> list[ColumnMeta]:
    """Build column metadata list from a DataFrame."""
    cols: list[ColumnMeta] = []
    for col_name in df.columns:
        cols.append(
            ColumnMeta(
                name=str(col_name),
                dtype=str(df[col_name].dtype),
                is_numeric=bool(pd.api.types.is_numeric_dtype(df[col_name])),
            )
        )
    return cols


# ---------------------------------------------------------------------------
# Async persistence helpers (fire-and-forget)
# ---------------------------------------------------------------------------

async def _persist_to_r2(
    r2_key: str,
    content: bytes,
    content_type: str,
) -> None:
    """Upload file bytes to R2 in a background task (non-blocking)."""
    try:
        # boto3 is synchronous — run in thread pool
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(
            None,
            r2.upload_file_bytes,
            r2_key,
            content,
            content_type,
        )
        logger.info("R2 upload complete: %s", r2_key)
    except Exception as exc:
        logger.error("R2 background upload failed: %s", exc)


async def _persist_metadata(
    clerk_user_id: str,
    file_id: str,
    file_name: str,
    r2_key: str,
) -> None:
    """Upsert user + create dataset row in Supabase (non-blocking)."""
    try:
        loop = asyncio.get_running_loop()
        # Upsert user
        db_user_id = await loop.run_in_executor(
            None, supa.upsert_user, clerk_user_id, None, None,
        )
        if db_user_id:
            # Create dataset
            await loop.run_in_executor(
                None, supa.create_dataset, db_user_id, file_name, r2_key, file_id,
            )
        logger.info("Supabase metadata persisted for file_id=%s", file_id)
    except Exception as exc:
        logger.error("Supabase background persist failed: %s", exc)


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@router.post("/upload/presign", response_model=PresignResponse)
async def upload_presign(body: PresignRequest, request: Request) -> PresignResponse:
    """Generate a presigned upload URL for direct R2 upload from the frontend.

    Requires ``x-clerk-user-id`` header.  Returns the presigned URL and R2 key.
    """
    clerk_user_id = get_current_user_id(request)
    if not clerk_user_id:
        raise HTTPException(
            status_code=401,
            detail="x-clerk-user-id header is required for presigned URL generation.",
        )

    if not r2.is_available():
        raise HTTPException(
            status_code=503,
            detail="R2 storage is not configured.",
        )

    dataset_id = str(uuid.uuid4())
    ext = _get_extension(body.file_name) or ".csv"
    r2_key = r2.make_dataset_key(clerk_user_id, dataset_id, ext)

    upload_url = r2.generate_presigned_upload_url(r2_key)
    if not upload_url:
        raise HTTPException(
            status_code=500,
            detail="Failed to generate presigned URL.",
        )

    return PresignResponse(upload_url=upload_url, r2_key=r2_key)


@router.post("/upload", response_model=UploadResponse)
async def upload(request: Request, files: list[UploadFile] = File(...)) -> UploadResponse:
    """Accept one or more files, parse data and context, return metadata.

    Phase 1.5+: supports up to 5 files per request, 20MB each.
    Multiple primary files (Excel/CSV) are allowed — the first one is used as the
    main dataset. Context files (.docx/.pptx) are combined.
    Also persists to R2 + Supabase asynchronously if available.
    """

    # Extract user identity (None for anonymous)
    clerk_user_id = get_current_user_id(request)

    # 0. File count check
    if len(files) > MAX_FILES_PER_UPLOAD:
        raise HTTPException(
            status_code=422,
            detail=f"Maximum {MAX_FILES_PER_UPLOAD} files per upload. Got {len(files)}.",
        )

    # 1. Read all files and validate extensions + sizes
    file_data: list[tuple[str, str, bytes]] = []  # (filename, ext, content)

    for f in files:
        filename = f.filename or "unknown"
        ext = _get_extension(filename)

        # Extension check
        if ext not in ALLOWED_ALL:
            raise HTTPException(
                status_code=415,
                detail=f"Unsupported file type '{ext}'. Allowed: .xlsx, .csv, .docx, .pptx",
            )

        content = await f.read()

        # Size check
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413,
                detail=f"File '{filename}' exceeds the 20 MB limit.",
            )

        file_data.append((filename, ext, content))

    # 2. Separate primary and context files
    primary_files = [(fn, ext, data) for fn, ext, data in file_data if ext in ALLOWED_PRIMARY]
    if len(primary_files) == 0:
        raise HTTPException(
            status_code=422,
            detail="No primary data file (.xlsx or .csv) provided. At least one is required.",
        )

    # 3. Parse the first primary file as the main dataset
    primary_fn, primary_ext, primary_content = primary_files[0]
    if primary_ext == ".xlsx":
        df = _parse_excel(primary_content)
    else:
        df = _parse_csv(primary_content)

    # Strip column name whitespace (but do NOT lowercase)
    df.columns = df.columns.str.strip()

    # 4. Parse context files
    context_parts: list[str] = []
    for fn, ext, data in file_data:
        if ext == ".docx":
            context_parts.append(_extract_docx_text(data))
        elif ext == ".pptx":
            context_parts.append(_extract_pptx_text(data))

    context_text: str | None = "\n".join(context_parts) if context_parts else None

    # 5. Build metadata
    columns = _detect_columns(df)
    row_count = len(df)
    file_id = str(uuid.uuid4())

    # 5b. Compute content hash (SHA-256) for deduplication
    import hashlib
    content_hash = hashlib.sha256(primary_content).hexdigest()

    # 6. Build R2 key (if user is authenticated)
    r2_key: str | None = None
    if clerk_user_id and r2.is_available():
        r2_key = r2.make_dataset_key(clerk_user_id, file_id, primary_ext)

    # 7. Store entry in-memory (cache)
    entry = FileEntry(
        file_id=file_id,
        dataframe=df,
        columns=[c.model_dump() for c in columns],
        row_count=row_count,
        context_text=context_text,
        r2_key=r2_key,
        user_id=clerk_user_id,
        content_hash=content_hash,
        file_name=primary_fn,
        file_type=primary_ext,
    )
    await set_entry(entry)

    # 8. Async persistence — fire-and-forget (non-blocking)
    if r2_key and r2.is_available():
        content_type = "text/csv" if primary_ext == ".csv" else "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        asyncio.create_task(_persist_to_r2(r2_key, primary_content, content_type))

    if clerk_user_id and r2_key and supa.is_available():
        asyncio.create_task(_persist_metadata(clerk_user_id, file_id, primary_fn, r2_key))

    return UploadResponse(
        file_id=file_id,
        columns=columns,
        row_count=row_count,
        context_extracted=context_text is not None,
    )

